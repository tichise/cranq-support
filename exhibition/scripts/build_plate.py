"""Build the Cranq exhibition plate (.pptx) for 生成AIなんでも展示会.

Generates a single A4-portrait slide that imports cleanly into Google Slides
via File → Open → Upload, with an embedded QR code linking to the support page.
"""

from __future__ import annotations

import io
from pathlib import Path

import qrcode
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Mm, Pt

ROOT = Path(__file__).resolve().parents[2]
EXHIBITION = ROOT / "exhibition"
OUT_PPTX = EXHIBITION / "cranq-plate.pptx"
OUT_QR = EXHIBITION / "qr-support.png"
APP_ICON = ROOT / "images" / "app-icon.jpg"

SUPPORT_URL = "https://tichise.github.io/cranq-support/"

# --- design tokens (synthesized from the four-role brief) -------------------
PAPER = RGBColor(0xFB, 0xF7, 0xF2)
INK = RGBColor(0x1A, 0x14, 0x10)
INK_SUB = RGBColor(0x5C, 0x4F, 0x44)
INK_MUTED = RGBColor(0x9C, 0x8E, 0x81)
ACCENT = RGBColor(0xF9, 0x73, 0x16)
ACCENT_DEEP = RGBColor(0xC2, 0x41, 0x0C)
SOFT = RGBColor(0xF4, 0xED, 0xE3)
DEEP = RGBColor(0xE8, 0xDE, 0xCF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

JP = "Noto Sans JP"
EN = "Inter"
MONO = "JetBrains Mono"


def make_qr(url: str, path: Path) -> Path:
    qr = qrcode.QRCode(
        version=None,
        error_correction=qrcode.constants.ERROR_CORRECT_Q,
        box_size=20,
        border=2,
    )
    qr.add_data(url)
    qr.make(fit=True)
    img = qr.make_image(fill_color="black", back_color="white").convert("RGB")
    img.save(path, format="PNG", optimize=True)
    return path


def set_solid_fill(shape, rgb: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def no_line(shape) -> None:
    shape.line.fill.background()


def line(shape, rgb: RGBColor, width_pt: float) -> None:
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width_pt)


def add_rect(slide, x_mm, y_mm, w_mm, h_mm, fill=None, border=None, border_width=0.4):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Mm(x_mm), Mm(y_mm), Mm(w_mm), Mm(h_mm)
    )
    if fill is None:
        shape.fill.background()
    else:
        set_solid_fill(shape, fill)
    if border is None:
        no_line(shape)
    else:
        line(shape, border, border_width)
    shape.shadow.inherit = False
    return shape


def add_text(
    slide,
    x_mm,
    y_mm,
    w_mm,
    h_mm,
    runs,
    *,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    line_spacing=1.4,
):
    """runs: list of dicts {text, font, size_pt, color, bold?, tracking_em?}.

    Multiple runs go on a single paragraph. Use "\n" inside text to break lines.
    """
    box = slide.shapes.add_textbox(Mm(x_mm), Mm(y_mm), Mm(w_mm), Mm(h_mm))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor

    paragraphs_specs: list[list[dict]] = [[]]
    for r in runs:
        text = r["text"]
        parts = text.split("\n")
        for i, part in enumerate(parts):
            if i > 0:
                paragraphs_specs.append([])
            paragraphs_specs[-1].append({**r, "text": part})

    for p_idx, p_runs in enumerate(paragraphs_specs):
        para = tf.paragraphs[0] if p_idx == 0 else tf.add_paragraph()
        para.alignment = align
        para.line_spacing = line_spacing
        for i, r in enumerate(p_runs):
            run = para.add_run() if i > 0 else para.add_run()
            run.text = r["text"]
            f = run.font
            f.name = r.get("font", JP)
            f.size = Pt(r["size_pt"])
            f.bold = r.get("bold", False)
            f.color.rgb = r.get("color", INK)
    return box


def main() -> None:
    EXHIBITION.mkdir(parents=True, exist_ok=True)
    make_qr(SUPPORT_URL, OUT_QR)

    prs = Presentation()
    prs.slide_width = Mm(210)
    prs.slide_height = Mm(297)

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Paper background
    add_rect(slide, 0, 0, 210, 297, fill=PAPER)

    M = 14  # outer margin (mm)
    CW = 210 - 2 * M  # content width = 182

    # ---- A. Header band -----------------------------------------------------
    # App icon (12 mm) + wordmark + exhibit tag
    icon_size = 12
    slide.shapes.add_picture(
        str(APP_ICON), Mm(M), Mm(M), Mm(icon_size), Mm(icon_size)
    )
    add_text(
        slide,
        M + icon_size + 3,
        M + 0.5,
        60,
        icon_size,
        [
            {"text": "Cranq", "font": "Inter Tight", "size_pt": 20, "bold": True},
            {"text": "  iOS · Cycling & Running",
             "font": EN, "size_pt": 9, "color": INK_MUTED},
        ],
        anchor=MSO_ANCHOR.MIDDLE,
        line_spacing=1.0,
    )
    add_text(
        slide,
        210 - M - 80,
        M + 0.5,
        80,
        icon_size,
        [
            {"text": "EXHIBIT  ·  ",
             "font": MONO, "size_pt": 8, "color": INK_MUTED, "bold": True},
            {"text": "#生成AIなんでも展示会",
             "font": JP, "size_pt": 9, "color": INK_SUB, "bold": True},
        ],
        align=PP_ALIGN.RIGHT,
        anchor=MSO_ANCHOR.MIDDLE,
        line_spacing=1.0,
    )

    # ---- accent bar under header -------------------------------------------
    add_rect(slide, M, M + icon_size + 4, 22, 0.8, fill=ACCENT)

    # ---- B. Hero -----------------------------------------------------------
    hero_y = M + icon_size + 10
    add_text(
        slide,
        M,
        hero_y,
        CW,
        8,
        [
            {"text": "EXHIBITION  /  生成AIなんでも展示会",
             "font": MONO, "size_pt": 9, "color": ACCENT_DEEP, "bold": True},
        ],
        line_spacing=1.0,
    )
    add_text(
        slide,
        M,
        hero_y + 8,
        CW,
        38,
        [
            {"text": "走った記録を、\n",
             "font": JP, "size_pt": 32, "color": INK, "bold": True},
            {"text": "AI", "font": "Inter Tight",
             "size_pt": 32, "color": ACCENT_DEEP, "bold": True},
            {"text": "と一緒に読み返す。",
             "font": JP, "size_pt": 32, "color": INK, "bold": True},
        ],
        line_spacing=1.35,
    )
    add_text(
        slide,
        M,
        hero_y + 30,
        CW,
        7,
        [
            {"text": "Your workouts, in your words — by your key.",
             "font": "Inter Tight", "size_pt": 13, "color": INK_SUB},
        ],
        line_spacing=1.2,
    )

    # ---- C. Lead -----------------------------------------------------------
    lead_y = hero_y + 42
    add_text(
        slide,
        M,
        lead_y,
        CW,
        20,
        [
            {"text": (
                "Cranq は、サイクリングとランニングのための iOS アプリ。"
                "HealthKit に記録されたワークアウトのルート・標高・心拍・天気を "
                "1 画面でまとめて振り返り、必要なら AI コーチに読ませて感想をもらえます。"
            ), "font": JP, "size_pt": 11, "color": INK_SUB},
        ],
        line_spacing=1.7,
    )

    # ---- D. AI section ------------------------------------------------------
    ai_y = lead_y + 22
    # section label
    add_text(
        slide,
        M,
        ai_y,
        CW,
        6,
        [
            {"text": "WHERE  AI  LIVES",
             "font": MONO, "size_pt": 10, "color": ACCENT_DEEP, "bold": True},
            {"text": "   AI はどこで動いているか",
             "font": JP, "size_pt": 11, "color": INK, "bold": True},
        ],
        line_spacing=1.0,
    )
    # 3-step flow
    flow_y = ai_y + 9
    flow_h = 22
    bubble_w = 50
    gap = (CW - 3 * bubble_w) / 2
    steps = [
        ("01", "あなたの\nOpenAI キーを設定"),
        ("02", "端末から OpenAI へ\n直接送信"),
        ("03", "ワークアウトを読んだ\n自然言語フィードバック"),
    ]
    for i, (num, label) in enumerate(steps):
        bx = M + i * (bubble_w + gap)
        # bubble
        b = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Mm(bx), Mm(flow_y), Mm(bubble_w), Mm(flow_h)
        )
        set_solid_fill(b, SOFT)
        line(b, DEEP, 0.5)
        b.shadow.inherit = False
        # number badge
        n = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Mm(bx + 4), Mm(flow_y + 4), Mm(7), Mm(7)
        )
        set_solid_fill(n, ACCENT)
        no_line(n)
        n.shadow.inherit = False
        n.text_frame.margin_left = n.text_frame.margin_right = 0
        n.text_frame.margin_top = n.text_frame.margin_bottom = 0
        n.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        para = n.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = num
        run.font.name = MONO
        run.font.size = Pt(8)
        run.font.bold = True
        run.font.color.rgb = WHITE
        # label
        add_text(
            slide,
            bx + 13,
            flow_y + 4,
            bubble_w - 16,
            flow_h - 6,
            [{"text": label, "font": JP, "size_pt": 9.5, "color": INK}],
            line_spacing=1.35,
        )
        # arrow
        if i < 2:
            ax = bx + bubble_w + 1
            ay = flow_y + flow_h / 2
            add_text(
                slide,
                ax,
                ay - 4,
                gap - 2,
                8,
                [{"text": "→", "font": EN, "size_pt": 16, "color": ACCENT, "bold": True}],
                align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE,
                line_spacing=1.0,
            )

    # BYOK pill (centered under flow)
    pill_w = 60
    pill_h = 8
    pill_x = M + (CW - pill_w) / 2
    pill_y = flow_y + flow_h + 3
    pill = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Mm(pill_x), Mm(pill_y), Mm(pill_w), Mm(pill_h)
    )
    set_solid_fill(pill, ACCENT)
    no_line(pill)
    pill.shadow.inherit = False
    pill.text_frame.margin_left = pill.text_frame.margin_right = 0
    pill.text_frame.margin_top = pill.text_frame.margin_bottom = 0
    pill.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    pp = pill.text_frame.paragraphs[0]
    pp.alignment = PP_ALIGN.CENTER
    r1 = pp.add_run()
    r1.text = "BYOK  ·  "
    r1.font.name = MONO
    r1.font.size = Pt(10)
    r1.font.bold = True
    r1.font.color.rgb = WHITE
    r2 = pp.add_run()
    r2.text = "中継サーバーなし。鍵はあなたのもの。"
    r2.font.name = JP
    r2.font.size = Pt(10)
    r2.font.bold = True
    r2.font.color.rgb = WHITE

    # ---- E. Effects (2x2) --------------------------------------------------
    eff_y = pill_y + pill_h + 7
    add_text(
        slide,
        M,
        eff_y,
        CW,
        6,
        [
            {"text": "WHY  IT  WORKS",
             "font": MONO, "size_pt": 10, "color": ACCENT_DEEP, "bold": True},
            {"text": "   設計上の効果",
             "font": JP, "size_pt": 11, "color": INK, "bold": True},
        ],
        line_spacing=1.0,
    )
    cards = [
        ("数値だけを渡す",
         "氏名や連絡先は送らない。AI に渡すのは HealthKit の運動指標。"),
        ("BYOK（自分の鍵）",
         "AI 利用ぶんはあなたの OpenAI 従量。サブスクや追加課金はゼロ。"),
        ("中継サーバーなし",
         "通信は端末から OpenAI へ直送。Cranq 側にログは残らない。"),
        ("AI はオプトイン",
         "鍵を設定しない限り AI は動かない。記録閲覧と統計は単独で完結。"),
    ]
    card_w = (CW - 4) / 2
    card_h = 19
    for i, (h, body) in enumerate(cards):
        col = i % 2
        row = i // 2
        cx = M + col * (card_w + 4)
        cy = eff_y + 9 + row * (card_h + 3)
        card = add_rect(slide, cx, cy, card_w, card_h, fill=SOFT, border=DEEP, border_width=0.5)
        # left accent stripe
        add_rect(slide, cx, cy, 1.2, card_h, fill=ACCENT)
        add_text(
            slide,
            cx + 4,
            cy + 2,
            card_w - 6,
            6,
            [{"text": h, "font": JP, "size_pt": 11, "color": INK, "bold": True}],
            line_spacing=1.0,
        )
        add_text(
            slide,
            cx + 4,
            cy + 8,
            card_w - 6,
            card_h - 9,
            [{"text": body, "font": JP, "size_pt": 9, "color": INK_SUB}],
            line_spacing=1.5,
        )

    # ---- F. Tagline strip --------------------------------------------------
    tag_y = eff_y + 9 + 2 * (card_h + 3) + 3
    add_text(
        slide,
        M,
        tag_y,
        CW,
        7,
        [
            {"text": "「計画する AI」ではなく、",
             "font": JP, "size_pt": 11, "color": INK_SUB},
            {"text": "振り返る AI",
             "font": JP, "size_pt": 11, "color": ACCENT_DEEP, "bold": True},
            {"text": "。日記を続けられない人の、書いてくれる版。",
             "font": JP, "size_pt": 11, "color": INK_SUB},
        ],
        line_spacing=1.4,
    )

    # ---- G. QR + footer ----------------------------------------------------
    qr_size = 32
    qr_x = M
    qr_y = 297 - M - qr_size
    # QR white card
    add_rect(slide, qr_x, qr_y, qr_size, qr_size, fill=WHITE, border=DEEP, border_width=0.5)
    pad = 1.5
    slide.shapes.add_picture(
        str(OUT_QR),
        Mm(qr_x + pad),
        Mm(qr_y + pad),
        Mm(qr_size - 2 * pad),
        Mm(qr_size - 2 * pad),
    )
    # caption right of QR
    cap_x = qr_x + qr_size + 5
    cap_w = CW - qr_size - 5
    add_text(
        slide,
        cap_x,
        qr_y + 1,
        cap_w,
        6,
        [{"text": "SUPPORT  ·  使い方とサポート",
          "font": MONO, "size_pt": 9, "color": ACCENT_DEEP, "bold": True}],
        line_spacing=1.0,
    )
    add_text(
        slide,
        cap_x,
        qr_y + 7,
        cap_w,
        10,
        [{"text": "tichise.github.io/\ncranq-support",
          "font": MONO, "size_pt": 12, "color": INK, "bold": True}],
        line_spacing=1.25,
    )
    add_text(
        slide,
        cap_x,
        qr_y + qr_size - 14,
        cap_w,
        6,
        [{"text": "App Store でも配信中（リンクはサポートページから）",
          "font": JP, "size_pt": 9, "color": INK_SUB}],
        line_spacing=1.2,
    )
    add_text(
        slide,
        cap_x,
        qr_y + qr_size - 7,
        cap_w,
        6,
        [
            {"text": "作: ", "font": JP, "size_pt": 9, "color": INK_MUTED},
            {"text": "tichise", "font": EN, "size_pt": 9, "color": INK, "bold": True},
            {"text": "  ·  @tichise",
             "font": MONO, "size_pt": 9, "color": INK_MUTED},
        ],
        line_spacing=1.2,
    )

    prs.save(OUT_PPTX)
    print(f"wrote {OUT_PPTX.relative_to(ROOT)}")
    print(f"wrote {OUT_QR.relative_to(ROOT)}")


if __name__ == "__main__":
    main()
